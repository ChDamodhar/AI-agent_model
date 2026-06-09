import os
import pandas as pd
import numpy as np
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import Dict, Any

class ReportAgent:
    def __init__(self):
        pass

    def generate_pdf_report(self, report_data: dict, reports_dir: str, file_id: str) -> str:
        pdf_filename = f"{file_id}_final_report.pdf"
        pdf_path = os.path.join(reports_dir, pdf_filename)
        
        doc = SimpleDocTemplate(pdf_path, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
        story = []
        
        styles = getSampleStyleSheet()
        
        # Premium design styles
        title_style = ParagraphStyle(
            'CoverTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=28,
            leading=34,
            textColor=colors.HexColor('#0F172A'),
            spaceAfter=15
        )
        
        subtitle_style = ParagraphStyle(
            'CoverSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=14,
            leading=18,
            textColor=colors.HexColor('#6366F1'),
            spaceAfter=30
        )
        
        h1_style = ParagraphStyle(
            'Heading1_Custom',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=18,
            leading=22,
            textColor=colors.HexColor('#1E293B'),
            spaceBefore=15,
            spaceAfter=10,
            keepWithNext=True
        )

        h2_style = ParagraphStyle(
            'Heading2_Custom',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=13,
            leading=16,
            textColor=colors.HexColor('#475569'),
            spaceBefore=10,
            spaceAfter=6,
            keepWithNext=True
        )

        body_style = ParagraphStyle(
            'Body_Custom',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#334155'),
            spaceAfter=8
        )
        
        meta_style = ParagraphStyle(
            'Meta_Custom',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=9,
            leading=12,
            textColor=colors.HexColor('#64748B'),
            spaceAfter=15
        )

        # 1. Title Section
        story.append(Paragraph("DataMind AI", title_style))
        story.append(Paragraph("Multi-Agent Autonomous Data Science Platform — Final Report", subtitle_style))
        story.append(Paragraph(f"Dataset Filename: <b>{report_data.get('filename', 'N/A')}</b>", body_style))
        story.append(Paragraph(f"File ID: <b>{file_id}</b>", meta_style))
        story.append(Spacer(1, 20))
        
        # 2. Data Cleaning Section
        story.append(Paragraph("1. Data Cleaning Summary", h1_style))
        clean_stats = report_data.get("cleaning_report", {})
        cleaning_text = (
            f"During the data cleaning stage, missing values, duplicates, and outliers were processed. "
            f"<b>Original dimensions:</b> {clean_stats.get('original_shape', ['N/A'])[0]} rows × {clean_stats.get('original_shape', ['N/A'])[1]} columns. "
            f"<b>Cleaned dimensions:</b> {clean_stats.get('cleaned_shape', ['N/A'])[0]} rows × {clean_stats.get('cleaned_shape', ['N/A'])[1]} columns.<br/>"
            f"• <b>Missing values imputed:</b> {clean_stats.get('missing_values', 0)} values.<br/>"
            f"• <b>Duplicate rows removed:</b> {clean_stats.get('duplicates_removed', 0)} rows.<br/>"
            f"• <b>Outliers capped:</b> {clean_stats.get('outliers_detected', 0)} outliers detected via IQR."
        )
        story.append(Paragraph(cleaning_text, body_style))
        story.append(Spacer(1, 15))
        
        # 3. EDA Section
        story.append(Paragraph("2. Exploratory Data Analysis (EDA)", h1_style))
        eda_report = report_data.get("eda_report", {})
        target_column = eda_report.get("target_column", "N/A")
        story.append(Paragraph(f"The detected target column for machine learning is: <b>{target_column}</b>.", body_style))
        
        plots = report_data.get("eda_plots", {})
        corr_plot_name = plots.get("correlation_heatmap")
        if corr_plot_name:
            corr_plot_path = os.path.join(reports_dir, corr_plot_name)
            if os.path.exists(corr_plot_path):
                story.append(Paragraph("<b>Correlation Heatmap:</b>", h2_style))
                story.append(Image(corr_plot_path, width=400, height=240))
                story.append(Spacer(1, 10))
                
        target_plot_name = plots.get("target_distribution")
        if target_plot_name:
            target_plot_path = os.path.join(reports_dir, target_plot_name)
            if os.path.exists(target_plot_path):
                story.append(Paragraph("<b>Target Column Distribution:</b>", h2_style))
                story.append(Image(target_plot_path, width=400, height=240))
                story.append(Spacer(1, 10))

        story.append(PageBreak())
        
        # 4. Feature Engineering Section
        story.append(Paragraph("3. Feature Engineering", h1_style))
        fe_report = report_data.get("feature_engineering_report", {})
        detected = fe_report.get("columns_detected", {})
        story.append(Paragraph(
            f"Columns were parsed and categorized: "
            f"<b>Numerical:</b> {len(detected.get('numerical', []))}, "
            f"<b>Categorical:</b> {len(detected.get('categorical', []))}, "
            f"<b>Datetime:</b> {len(detected.get('datetime', []))}, "
            f"<b>Text:</b> {len(detected.get('text', []))}.",
            body_style
        ))
        
        story.append(Paragraph("<b>Transformations Applied:</b>", h2_style))
        for step in fe_report.get("transformations_applied", []):
            story.append(Paragraph(f"• {step}", body_style))
        story.append(Spacer(1, 15))
        
        # 5. Model Selection Leaderboard
        story.append(Paragraph("4. Model Selection & Tuning", h1_style))
        model_report = report_data.get("model_selection_report", {})
        best_model_name = model_report.get("best_model_name", "N/A")
        problem_type = model_report.get("problem_type", "classification")
        
        story.append(Paragraph(f"Problem Type: <b>{problem_type}</b>. The best model selected is: <b>{best_model_name}</b>.", body_style))
        
        leaderboard = model_report.get("leaderboard", [])
        if leaderboard:
            story.append(Paragraph("<b>Leaderboard:</b>", h2_style))
            if problem_type == "classification":
                table_data = [["Model Name", "Accuracy", "F1 Score", "ROC-AUC"]]
                for row in leaderboard:
                    table_data.append([
                        row.get("model_name"),
                        f"{row.get('accuracy', 0.0):.4f}",
                        f"{row.get('f1_score', 0.0):.4f}",
                        f"{row.get('roc_auc', 0.5):.4f}"
                    ])
            else:
                table_data = [["Model Name", "R2 Score", "MSE", "MAE"]]
                for row in leaderboard:
                    table_data.append([
                        row.get("model_name"),
                        f"{row.get('r2_score', 0.0):.4f}",
                        f"{row.get('mse', 0.0):.4f}",
                        f"{row.get('mae', 0.0):.4f}"
                    ])
                    
            table = Table(table_data, colWidths=[180, 100, 100, 100])
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0F172A')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#E2E8F0')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#FFFFFF'), colors.HexColor('#F8FAFC')]),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
            ]))
            story.append(table)
            story.append(Spacer(1, 15))
            
        tuning_report = report_data.get("tuning_report", {})
        if tuning_report:
            story.append(Paragraph("<b>Hyperparameter Tuning (Optuna):</b>", h2_style))
            tunable = tuning_report.get("tunable", True)
            if tunable:
                tune_text = (
                    f"Optuna optimized the hyperparameters of the {tuning_report.get('model_name')}.<br/>"
                    f"• <b>Baseline Metric:</b> {tuning_report.get('baseline_metric', 0.0):.4f}<br/>"
                    f"• <b>Tuned Metric:</b> {tuning_report.get('tuned_metric', 0.0):.4f}<br/>"
                    f"• <b>Best Parameters:</b> {tuning_report.get('best_params', {})}"
                )
            else:
                tune_text = "Model is not tunable (e.g. LinearRegression)."
            story.append(Paragraph(tune_text, body_style))
            story.append(Spacer(1, 15))

        # 7. Business Insights Section (Phase 9)
        insights = report_data.get("business_insights", [])
        if insights:
            story.append(Paragraph("6. Business Insights", h1_style))
            for insight in insights:
                story.append(Paragraph(f"• {insight}", body_style))
            story.append(Spacer(1, 15))
        
        # Build the PDF
        doc.build(story)
        return pdf_filename

    def generate_pptx_report(self, report_data: dict, reports_dir: str, file_id: str) -> str:
        pptx_filename = f"{file_id}_final_presentation.pptx"
        pptx_path = os.path.join(reports_dir, pptx_filename)
        
        prs = Presentation()
        
        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "DataMind AI Platform Insights"
        subtitle.text = f"Autonomous ML Analysis for {report_data.get('filename', 'Dataset')}\nFile ID: {file_id}"
        
        # Slide 2: Data Cleaning
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Data Cleaning Summary"
        content = slide.placeholders[1]
        clean_stats = report_data.get("cleaning_report", {})
        content.text = (
            f"• Initial Data Dimensions: {clean_stats.get('original_shape', ['N/A'])[0]} rows x {clean_stats.get('original_shape', ['N/A'])[1]} columns\n"
            f"• Cleaned Data Dimensions: {clean_stats.get('cleaned_shape', ['N/A'])[0]} rows x {clean_stats.get('cleaned_shape', ['N/A'])[1]} columns\n"
            f"• Duplicate rows dropped: {clean_stats.get('duplicates_removed', 0)}\n"
            f"• Missing values imputed: {clean_stats.get('missing_values', 0)}\n"
            f"• Outliers capped via IQR: {clean_stats.get('outliers_detected', 0)} columns affected"
        )
        
        # Slide 3: EDA Plots
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Exploratory Data Analysis"
        plots = report_data.get("eda_plots", {})
        corr_plot_name = plots.get("correlation_heatmap")
        if corr_plot_name:
            corr_plot_path = os.path.join(reports_dir, corr_plot_name)
            if os.path.exists(corr_plot_path):
                slide.shapes.add_picture(corr_plot_path, Inches(0.5), Inches(1.8), width=Inches(4.2), height=Inches(3.2))
                
        target_plot_name = plots.get("target_distribution")
        if target_plot_name:
            target_plot_path = os.path.join(reports_dir, target_plot_name)
            if os.path.exists(target_plot_path):
                slide.shapes.add_picture(target_plot_path, Inches(5.0), Inches(1.8), width=Inches(4.2), height=Inches(3.2))
                
        # Slide 4: Feature Engineering
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Feature Engineering Details"
        content = slide.placeholders[1]
        fe_report = report_data.get("feature_engineering_report", {})
        detected = fe_report.get("columns_detected", {})
        
        bullets = [
            f"Column classification: {len(detected.get('numerical', []))} Numerical, {len(detected.get('categorical', []))} Categorical, {len(detected.get('datetime', []))} Datetime.",
            "Applied transformations:"
        ]
        for step in fe_report.get("transformations_applied", [])[:4]:
            bullets.append(f"  - {step}")
            
        content.text = "\n".join(bullets)
        
        # Slide 5: Model Selection & Leaderboard
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Model Selection Leaderboard"
        content = slide.placeholders[1]
        model_report = report_data.get("model_selection_report", {})
        best_model_name = model_report.get("best_model_name", "N/A")
        
        bullets = [
            f"• Best Model: {best_model_name}",
            f"• Problem Type: {model_report.get('problem_type', 'classification')}",
            "• Performance Leaderboard:"
        ]
        leaderboard = model_report.get("leaderboard", [])
        for i, row in enumerate(leaderboard[:3]):
            bullets.append(f"  {i+1}. {row.get('model_name')} (Score: {row.get('f1_score', row.get('r2_score', 0.0)):.4f})")
            
        content.text = "\n".join(bullets)
        
        # Slide 6: Explainability (SHAP Summary)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Model Explainability (SHAP)"
        explanation = report_data.get("explanation_report", {})
        if explanation:
            shap_plot_name = explanation.get("shap_plot_path")
            if shap_plot_name:
                shap_plot_path = os.path.join(reports_dir, shap_plot_name)
                if os.path.exists(shap_plot_path):
                    slide.shapes.add_picture(shap_plot_path, Inches(1.5), Inches(1.5), width=Inches(7.0), height=Inches(5.0))
                    
        prs.save(pptx_path)
        return pptx_filename
